# Auto-generate PPT (and optional PDF) for CommandDispatcher commands + keyboard shortcuts

import ast
from datetime import date
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import os
import comtypes.client
import textwrap

# -------------------------------------------------------------------
# CONFIG
# -------------------------------------------------------------------
EXPORT_TO_PDF = False          # also export to PDF via PowerPoint
FONT_SIZE = 12                 # content font size

COMMAND_DISPATCHER_PATH = r"c:\WC\HotSystem\CommandDispatcher.py"
APPLICATION_PATH = r"c:\WC\HotSystem\Application.py"

OUT_PPTX = "CommandDispatcher_Commands.pptx"
OUT_PDF = "CommandDispatcher_Commands.pdf"


# -------------------------------------------------------------------
# Helpers
# -------------------------------------------------------------------
def wrap_lines(lines, max_chars=70):
    """Soft-wrap lines to avoid overflowing the text boxes."""
    wrapped = []
    for line in lines:
        wrapped.extend(
            textwrap.wrap(
                line,
                width=max_chars,
                break_long_words=False,
                break_on_hyphens=False,
            )
        )
    return wrapped


def pretty_key_name(name: str) -> str:
    """Convert KeyboardKeys-style enum names to user-friendly key labels."""
    mapping = {
        "OEM_4": "[",
        "OEM_5": "\\",
        "OEM_6": "]",
        "OEM_PLUS": "=",
        "OEM_MINUS": "-",
        "SPACE_KEY": "Space",
        "UP_KEY": "Up Arrow",
        "DOWN_KEY": "Down Arrow",
        "LEFT_KEY": "Left Arrow",
        "RIGHT_KEY": "Right Arrow",
        "PAGEUP_KEY": "Page Up",
        "PAGEDOWN_KEY": "Page Down",
        "INSERT_KEY": "Insert",
        "DEL_KEY": "Delete",
        "HOME_KEY": "Home",
        "END_KEY": "End",
        "BACK_KEY": "Backspace",
        "ENTER_KEY": "Enter",
        "C_KEY": "C",
        "V_KEY": "V",
        "L_KEY": "L",
        "R_KEY": "R",
        "S_KEY": "S",
        "K_KEY": "K",
        "Q_KEY": "Q",
        "W_KEY": "W",
        "E_KEY": "E",
        "G_KEY": "G",
        "H_KEY": "H",
        "B_KEY": "B",
        "N_KEY": "N",
        "X_KEY": "X",
        "F_KEY": "F",
        "M_KEY": "M",
        "A_KEY": "A",
        "O_KEY": "O",
        "P_KEY": "P",
        "T_KEY": "T",
        "Y_KEY": "Y",
        "Z_KEY": "Z",
        "OEM_3": "`",
    }

    if name in mapping:
        return mapping[name]

    # KEY_0..KEY_9
    if name.startswith("KEY_") and len(name) == 5 and name[-1].isdigit():
        return name[-1]

    # A_KEY style
    if name.endswith("_KEY") and len(name) == 5 and name[0].isalpha():
        return name[0].upper()

    return name


# -------------------------------------------------------------------
# Highlighting console commands in green (docstring slides)
# -------------------------------------------------------------------
def build_sorted_commands(source_code: str):
    """Parse CommandDispatcher.py and return (handlers_mapping, sorted_cmds, func_docs)."""
    tree = ast.parse(source_code)
    class_node = next(
        (
            node
            for node in tree.body
            if isinstance(node, ast.ClassDef) and node.name == "CommandDispatcher"
        ),
        None,
    )
    if not class_node:
        raise ValueError("CommandDispatcher class not found.")

    # Extract handlers mapping from __init__
    handlers_mapping = {}
    for node in class_node.body:
        if isinstance(node, ast.FunctionDef) and node.name == "__init__":
            for sub in node.body:
                if isinstance(sub, ast.Assign):
                    for target in sub.targets:
                        if (
                            isinstance(target, ast.Attribute)
                            and isinstance(sub.value, ast.Dict)
                            and isinstance(target.value, ast.Name)
                            and target.value.id == "self"
                            and target.attr == "handlers"
                        ):
                            for key_node, val_node in zip(
                                sub.value.keys, sub.value.values
                            ):
                                if (
                                    isinstance(key_node, ast.Constant)
                                    and isinstance(key_node.value, str)
                                ):
                                    cmd_name = key_node.value
                                else:
                                    continue
                                handler_name = None
                                if isinstance(val_node, ast.Attribute):
                                    if (
                                        isinstance(val_node.value, ast.Name)
                                        and val_node.value.id == "self"
                                    ):
                                        handler_name = val_node.attr
                                elif isinstance(val_node, ast.Lambda):
                                    if isinstance(val_node.body, ast.Call):
                                        func = val_node.body.func
                                        if isinstance(func, ast.Attribute):
                                            if (
                                                isinstance(func.value, ast.Name)
                                                and func.value.id == "self"
                                            ):
                                                handler_name = func.attr
                                handlers_mapping[cmd_name] = handler_name
            break

    # Docstrings by function
    func_docs = {}
    for node in class_node.body:
        if isinstance(node, ast.FunctionDef):
            func_docs[node.name] = ast.get_docstring(node, clean=True)

    sorted_cmds = sorted(handlers_mapping.keys())
    return handlers_mapping, sorted_cmds, func_docs


def highlight_command_line(p, line, sorted_cmds):
    """
    Highlight console command prefixes (from CommandDispatcher) in green
    on command docstring slides.
    """
    line_stripped = line.strip()
    for cmd in sorted(sorted_cmds, key=lambda x: -len(x)):
        if line_stripped.startswith(cmd + " ") or line_stripped == cmd:
            cmd_idx = line.find(cmd)
            before = line[:cmd_idx]
            after = line[cmd_idx + len(cmd) :]

            r1 = p.add_run()
            r1.text = before
            r1.font.size = Pt(FONT_SIZE)

            r2 = p.add_run()
            r2.text = cmd
            r2.font.size = Pt(FONT_SIZE)
            r2.font.color.rgb = RGBColor(0, 128, 0)

            r3 = p.add_run()
            r3.text = after
            r3.font.size = Pt(FONT_SIZE)
            return

    # no known command prefix -> plain black
    p.text = line
    for r in p.runs:
        r.font.size = Pt(FONT_SIZE)


# -------------------------------------------------------------------
# Extract keyboard shortcuts from Application.py
# -------------------------------------------------------------------
def first_sentence_from_doc(doc: str) -> str:
    if not doc:
        return ""
    return doc.strip().splitlines()[0].strip()


def extract_shortcuts_from_application(path: str):
    """
    Parse Application.py (PyGuiOverlay class) and extract keyboard shortcuts:
    - Ctrl + key combinations from ctrl_actions dict
    - Shift + key combinations from shift_actions dict
    - Core history/editing shortcuts (manual, because they are if/elif blocks)
    Returns a list of (category, combo, description).
    """
    if not os.path.exists(path):
        return []

    with open(path, "r", encoding="utf-8") as f:
        src = f.read()

    try:
        tree = ast.parse(src)
    except SyntaxError:
        return []

    # Find PyGuiOverlay class
    class_node = None
    for node in tree.body:
        if isinstance(node, ast.ClassDef) and node.name == "PyGuiOverlay":
            class_node = node
            break
    if not class_node:
        return []

    methods = {
        n.name: n for n in class_node.body if isinstance(n, ast.FunctionDef)
    }

    def describe_method(name: str) -> str:
        fn = methods.get(name)
        if fn:
            doc = ast.get_docstring(fn, clean=True)
            if doc:
                return first_sentence_from_doc(doc)
        # fallback: from name
        base = name.lstrip("_")
        return base.replace("_", " ").capitalize()

    def extract_dict_shortcuts(
        func_name: str, dict_var: str, modifier_label: str
    ):
        out = []
        fn = methods.get(func_name)
        if not fn:
            return out
        for stmt in fn.body:
            if isinstance(stmt, ast.Assign) and any(
                isinstance(t, ast.Name) and t.id == dict_var
                for t in stmt.targets
            ):
                d = stmt.value
                if not isinstance(d, ast.Dict):
                    continue
                for key_node, val_node in zip(d.keys, d.values):
                    if not isinstance(key_node, ast.Attribute):
                        continue
                    if not isinstance(key_node.value, ast.Name):
                        continue
                    if key_node.value.id != "KeyboardKeys":
                        continue
                    key_name = key_node.attr
                    key_pretty = pretty_key_name(key_name)
                    combo = f"{modifier_label} + {key_pretty}"

                    desc = ""

                    # method attribute, e.g. self._increase_exposure
                    if isinstance(val_node, ast.Attribute):
                        desc = describe_method(val_node.attr)

                    # Name (might be local helper)
                    elif isinstance(val_node, ast.Name):
                        desc = describe_method(val_node.id)

                    # Lambda
                    elif isinstance(val_node, ast.Lambda):
                        body = val_node.body
                        if isinstance(body, ast.Call):
                            func = body.func
                            # run("...") calls
                            if (
                                isinstance(func, ast.Name)
                                and func.id == "run"
                                and body.args
                                and isinstance(body.args[0], ast.Constant)
                                and isinstance(body.args[0].value, str)
                            ):
                                cmd_str = body.args[0].value
                                desc = f'Run console command "{cmd_str}"'
                            # helper like _nudge_carrier(...)
                            elif isinstance(func, ast.Name):
                                desc = describe_method(func.id)

                    if not desc:
                        desc = "Perform shortcut action"

                    out.append((modifier_label, combo, desc))
        return out

    shortcuts = []

    # Ctrl + key
    shortcuts.extend(
        extract_dict_shortcuts(
            "handle_ctrl_shift_commands", "ctrl_actions", "Ctrl"
        )
    )

    # Shift + key
    shortcuts.extend(
        extract_dict_shortcuts(
            "handle_ctrl_shift_commands", "shift_actions", "Shift"
        )
    )

    # Manual additions from keyboard_callback for history / clipboard / focus
    shortcuts.extend(
        [
            (
                "History",
                "Up Arrow",
                "Go to previous command in console history.",
            ),
            (
                "History",
                "Down Arrow",
                "Go to next command in console history (or clear when at end).",
            ),
            (
                "Clipboard/Run",
                "` (backtick)",
                "Paste clipboard text into the command box and execute it.",
            ),
            (
                "Editing",
                "Backspace",
                "Delete last character from cmd_input (and focus command box).",
            ),
            (
                "Editing",
                "Shift + Backspace",
                "Clear the entire command input line.",
            ),
            (
                "Focus",
                "C / Space / Enter",
                "Focus the command input field (cmd_input).",
            ),
        ]
    )

    # Summarized Smaract / Picomotor movement descriptions
    shortcuts.extend(
        [
            (
                "Smaract (movement)",
                "Ctrl + Arrows/PageUp/PageDown",
                "Move Smaract stage in coarse steps along X/Y/Z axes.",
            ),
            (
                "Smaract (movement)",
                "Shift + Arrows/PageUp/PageDown",
                "Move Smaract stage in fine steps along X/Y/Z axes.",
            ),
            (
                "Smaract",
                "Ctrl/Shift + Space",
                "Log current Smaract position as a point.",
            ),
            (
                "Picomotor",
                "Alt + Arrows/PageUp/PageDown",
                "Move Picomotor axes in steps along X/Y/Z.",
            ),
        ]
    )

    return shortcuts


# -------------------------------------------------------------------
# MAIN GENERATION
# -------------------------------------------------------------------
# 1. Load CommandDispatcher source and parse
with open(COMMAND_DISPATCHER_PATH, "r", encoding="utf-8") as f:
    cmd_source = f.read()

handlers_mapping, sorted_cmds, func_docs = build_sorted_commands(cmd_source)

# 2. Prepare PowerPoint
prs = Presentation()
layout = prs.slide_layouts[1]  # Title + Content
command_slide_refs = {}

# 3. Title slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title_slide.shapes.title.text = "Auto-Generated Help for Commands in Console GUI"
title_body = title_slide.placeholders[1].text_frame
title_body.text = (
    "This file was generated by the Python file 'generate_help_ppt.py'\n"
    f"Date: {date.today()}"
)

# 4. Index slide
index_slide = prs.slides.add_slide(prs.slide_layouts[5])
index_slide.shapes.title.text = "Command Index"

# 5. Build command slides
for cmd in sorted(handlers_mapping.keys()):
    handler = handlers_mapping[cmd]
    doc = func_docs.get(handler)
    lines = [line.strip() for line in doc.strip().splitlines() if line.strip()] if doc else []
    lines = wrap_lines(lines)
    is_long = len(lines) > 10

    slide = prs.slides.add_slide(prs.slide_layouts[5] if is_long else layout)
    slide.shapes.title.text = cmd

    if is_long:
        # 2 columns
        col_gap = Inches(0.4)
        col_width = (prs.slide_width - col_gap - Inches(1.0)) / 2
        height = prs.slide_height - Inches(1.5)
        left_col = Inches(0.5)
        right_col = left_col + col_width + col_gap
        top = Inches(1.0)

        mid = len(lines) // 2
        left_lines = lines[:mid]
        right_lines = lines[mid:]

        left_box = slide.shapes.add_textbox(left_col, top, col_width, height).text_frame
        right_box = slide.shapes.add_textbox(right_col, top, col_width, height).text_frame

        for box, chunk in ((left_box, left_lines), (right_box, right_lines)):
            for line in chunk:
                p = box.add_paragraph()
                highlight_command_line(p, line, sorted_cmds)

        tf_for_usage = right_box
    else:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for line in lines:
            p = tf.add_paragraph()
            highlight_command_line(p, line, sorted_cmds)
        tf_for_usage = tf

    # Usage example line
    p = tf_for_usage.add_paragraph()
    p.text = "Usage Example:"
    p.runs[0].font.size = Pt(FONT_SIZE)
    p.runs[0].font.bold = True
    usage = f"{cmd} [arg]" if cmd not in {"exit", "help", "quit"} else cmd
    p2 = tf_for_usage.add_paragraph()
    highlight_command_line(p2, usage, sorted_cmds)

    command_slide_refs[cmd] = slide

# 6. Fill command index (5 columns, each command linked to its slide)
num_columns = 5
col_width = prs.slide_width / num_columns - Inches(0.2)
row_height = Inches(0.4)
margin_top = Inches(1.2)
margin_left = Inches(0.2)

commands_sorted = sorted(handlers_mapping.keys())

for idx, cmd in enumerate(commands_sorted):
    col = idx % num_columns
    row = idx // num_columns
    left = margin_left + col * (col_width + Inches(0.2))
    top = margin_top + row * row_height
    width = col_width
    height = row_height

    box = index_slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = ""
    run = p.add_run()
    run.text = cmd
    run.font.size = Pt(FONT_SIZE)
    run.font.color.rgb = RGBColor(0, 128, 0)
    run.hyperlink.target_slide = command_slide_refs[cmd]

# 7. History slide (1,2,3,1?,2?,3?, ?)
history_slide = prs.slides.add_slide(layout)
history_slide.shapes.title.text = "Command History Access"
tf = history_slide.placeholders[1].text_frame
tf.clear()
history_lines = [
    "1  - Re-run the most recent command",
    "2  - Re-run the 2nd most recent command",
    "3  - Re-run the 3rd most recent command",
    "1? - Show the most recent command without executing",
    "2? - Show the 2nd most recent command without executing",
    "3? - Show the 3rd most recent command without executing",
    "?  - List all recent command history",
]
for line in history_lines:
    p = tf.add_paragraph()
    p.text = line
    for r in p.runs:
        r.font.size = Pt(FONT_SIZE)

# 8. Keyboard shortcuts slides from Application.py
shortcuts = extract_shortcuts_from_application(APPLICATION_PATH)

if shortcuts:
    # Sort a bit: by category then by combo
    shortcuts.sort(key=lambda x: (x[0], x[1]))

    max_per_slide = 18
    total = len(shortcuts)
    num_pages = (total + max_per_slide - 1) // max_per_slide

    for page_idx in range(num_pages):
        start = page_idx * max_per_slide
        end = min(start + max_per_slide, total)
        subset = shortcuts[start:end]

        slide = prs.slides.add_slide(layout)
        title = "Keyboard Shortcuts"
        if num_pages > 1:
            title += f" (page {page_idx + 1})"
        slide.shapes.title.text = title

        tf = slide.placeholders[1].text_frame
        tf.clear()

        for category, combo, desc in subset:
            p = tf.add_paragraph()
            p.text = ""
            # combo in green, description in black
            r1 = p.add_run()
            r1.text = f"{combo}: "
            r1.font.size = Pt(FONT_SIZE)
            r1.font.color.rgb = RGBColor(0, 128, 0)
            r1.font.bold = True

            r2 = p.add_run()
            r2.text = f"{desc}  [{category}]"
            r2.font.size = Pt(FONT_SIZE)

# 9. Save PPTX
prs.save(OUT_PPTX)

# 10. Optional PDF export through PowerPoint
if EXPORT_TO_PDF:
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    ppt = powerpoint.Presentations.Open(os.path.abspath(OUT_PPTX), WithWindow=False)
    ppt.SaveAs(os.path.abspath(OUT_PDF), FileFormat=32)  # 32 = PDF
    ppt.Close()
    powerpoint.Quit()
